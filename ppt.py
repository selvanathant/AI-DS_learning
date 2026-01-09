from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont
import os

# Create presentation
prs = Presentation()

# Ensure folder for icons exists
icon_folder = "icons"
os.makedirs(icon_folder, exist_ok=True)

# ---------- Functions ----------

def create_placeholder_icon(filename, color, letter):
    """Create a 100x100 px icon with a colored circle and a letter inside."""
    size = (100, 100)
    img = Image.new("RGBA", size, (255, 255, 255, 0))
    draw = ImageDraw.Draw(img)
    
    # Draw circle
    draw.ellipse((10, 10, 90, 90), fill=color)
    
    # Draw letter
    try:
        font = ImageFont.truetype("arial.ttf", 40)
    except:
        font = ImageFont.load_default()
    
    # Compute text size for centering
    w, h = draw.textbbox((0, 0), letter, font=font)[2:]
    draw.text(((size[0]-w)/2, (size[1]-h)/2), letter, fill="white", font=font)
    
    img.save(filename)
    return filename

def add_hex_background(slide):
    """Add hexagon-style techy background (just a gradient rectangle here)."""
    from pptx.shapes.autoshape import Shape
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls
    
    xml = f'''
    <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <p:nvSpPr>
            <p:cNvPr id="1" name="Background"/>
            <p:cNvSpPr/>
            <p:nvPr/>
        </p:nvSpPr>
        <p:spPr>
            <a:solidFill>
                <a:schemeClr val="accent1"/>
            </a:solidFill>
        </p:spPr>
    </p:sp>
    '''
    # This is just a placeholder rectangle. PPTX doesn't allow easy complex gradients.
    # Hexagon pattern can be drawn using shapes, but for simplicity, we use colored slide.
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)

def add_slide(title, points, icon_files):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Hex background
    add_hex_background(slide)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), prs.slide_width-1, Inches(1))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Add points with icons
    left = Inches(0.5)
    top = Inches(1.5)
    height = Inches(1)
    spacing = Inches(1.2)
    
    for i, point in enumerate(points):
        # White box
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + i*spacing, prs.slide_width-Inches(1), height)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        box.line.color.rgb = RGBColor(0, 0, 0)
        
        # Add point text
        tf = box.text_frame
        tf.text = point
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80)
        
        # Add icon on the left of box
        if i < len(icon_files):
            slide.shapes.add_picture(icon_files[i], left=Inches(0.2), top=top + i*spacing, width=Inches(0.8), height=Inches(0.8))

# ---------- Define slides ----------

slide_titles = [
    "Hello! I'm ChatGPT 🤖",
    "I Love Learning",
    "I Can Solve Problems",
    "I Read Books",
    "I Explore the World",
    "I Help You Code",
    "I Tell Stories",
    "I Love Art & Creativity",
    "I Keep Time",
    "I Make You Smile 😄"
]

slide_points = [
    ["Your friendly AI assistant", "I know a lot!", "I can chat endlessly"],
    ["Curious and curious again", "Ask me anything"],
    ["Math, science, logic", "Puzzles and challenges"],
    ["Fiction, non-fiction, tech books", "I love learning new stuff"],
    ["Science, geography, space", "Always exploring!"],
    ["Python, Java, C++", "Debugging & tips"],
    ["Short stories, poems, ideas", "Fun learning"],
    ["Drawings, design, colors", "Tech + fun"],
    ["Reminders, schedules, alarms", "Organized"],
    ["I try to make you happy", "Fun facts & jokes"]
]

# ---------- Create icons ----------
colors = [
    (255, 102, 102), (102, 204, 255), (255, 204, 102),
    (102, 255, 178), (255, 102, 255), (178, 102, 255),
    (255, 178, 102), (102, 178, 255), (178, 255, 102), (255, 102, 178)
]

icon_files = []
for i in range(10):
    file_path = os.path.join(icon_folder, f"icon_{i}.png")
    icon_files.append(create_placeholder_icon(file_path, colors[i], str(i+1)))

# ---------- Add slides ----------
for i in range(10):
    add_slide(slide_titles[i], slide_points[i], [icon_files[i]])

# ---------- Save PPT ----------
prs.save("ChatGPT_Bio_Techy_Figurines.pptx")
print("✅ 10-slide techy, child-friendly PPT created: ChatGPT_Bio_Techy_Figurines.pptx")
